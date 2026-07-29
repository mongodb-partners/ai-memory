#!/usr/bin/env python3
"""Build the Ai4 2026 booth-talk deck.

Rather than laying out slides from scratch, this rewrites the text of an existing
styled deck. Every slide in the base deck shares one layout — ``TextBox 4`` is the
title and ``TextBox 7`` is the body — and every body paragraph carries an identical
``<a:pPr>``. So we can swap the words and inherit the MongoDB brand styling, spacing,
and dark surface for free.

The base deck is the earlier draft of this same talk. Its slides 9 and 10 were
speaker-planning slides ("use the MCP-native assistant story", "four quick beats");
those slots become the demo and the episodic-memory slide. Its CTA pointed at
``ai-memory`` and ``memory-mcp``, both superseded by this package.

Usage::

    uv run python docs/talks/ai4-2026/build_deck.py
    uv run python docs/talks/ai4-2026/build_deck.py --base other.pptx --out mine.pptx

Content lives in ``SLIDES`` below and mirrors ``slides.md``. Speaker notes are written
into each slide's notes page, so ``slides.md`` stays the readable source and the deck
stays presentable on its own.
"""

from __future__ import annotations

import argparse
import copy
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ModuleNotFoundError:  # pragma: no cover - dependency hint
    sys.exit("python-pptx is required:  uv pip install python-pptx")

HERE = Path(__file__).resolve().parent
DEFAULT_BASE = Path.home() / "Downloads" / "Give_Your_Agents_a_Memory_Ai4.pptx"
DEFAULT_OUT = Path.home() / "Downloads" / "Give_Your_Agents_a_Memory_Ai4_v2.pptx"

TITLE_SHAPE = "TextBox 4"
BODY_SHAPE = "TextBox 7"

# Slide 1 is the title slide and has a different shape set; leave its text alone
# except for the subtitle, which is keyed by shape name.
TITLE_SLIDE_SUBTITLE = "TextBox 6"
TITLE_SLIDE_SUBTITLE_TEXT = (
    "Short-term state, long-term semantic memory, episodic recall, and hybrid "
    "retrieval — in one database."
)
TITLE_SLIDE_NOTES = (
    "0:20. Name, role, then straight in. Do not introduce yourself twice; booth "
    "crowds are already deciding whether to stay. 'Fifteen minutes, one idea, and a "
    "demo.'"
)

# (slide number, title, [body paragraphs], speaker notes)
SLIDES: list[tuple[int, str, list[str], str]] = [
    (
        2,
        "Your agent forgot. A bigger context window will not fix it.",
        [
            "It re-asks a question the user answered yesterday.",
            "It loses a decision that was made thirty turns ago.",
            "It has no idea what it actually did last Tuesday.",
            "Prompt engineering became context engineering for a reason: the "
            "bottleneck moved from how you word it to what data you put in front of "
            "the model.",
            "Bigger windows lose on three axes at once — cost is linear in tokens, "
            "attention degrades in the middle, and the window dies when the session "
            "does.",
        ],
        "1:10. Get the nods first — this is the cheapest audience buy-in in the talk. "
        "Land the pivot line: the moment the problem is 'what data, selected how, for "
        "which user', it is a database problem. Say the three-axis argument in one "
        "breath and move on. Do NOT get drawn into a needle-in-a-haystack benchmark "
        "debate at a booth.",
    ),
    (
        3,
        "Four kinds of memory. Most teams build one.",
        # Deliberately four short lines, not four explanations. This slide is read,
        # not narrated — the tiers get taught concretely on slides 5 to 7, and every
        # sentence added here is a sentence delivered twice.
        [
            "Short-term state — what are we talking about right now?",
            "Long-term semantic — what do I know about this user?",
            "Episodic — what did I actually do, and when?",
            "Semantic cache — have I answered this already?",
            "Most teams build one. You need four.",
        ],
        "1:15. Name the four, land 'most teams build one, you need four', and move. Do "
        "NOT explain reinforcement, merging, decay, or promotion here — each one is a "
        "question you would rather answer at the booth afterward than spend twenty "
        "seconds on now. Only if the crowd is leaning in, add ONE sentence: long-term "
        "stores distilled facts with scores, not transcripts; episodic is the one "
        "nobody builds, see slide 10; the cache makes the second identical question "
        "free. Booth audiences read faster than you talk — point at the slide and let "
        "them.",
    ),
    (
        4,
        "Same model. Same prompt. One difference: where its memory lives.",
        [
            "There is one toggle in the UI you are about to see.",
            "It does not change the model, the temperature, or the system prompt.",
            "It decides whether the agent is allowed to read and write memory.",
            "Everything else is held constant, on purpose — so the difference you see "
            "has exactly one cause.",
        ],
        "0:30. Set the frame BEFORE the demo. Otherwise half the crowd assumes you "
        "swapped models and the demo proves nothing. Say the words 'one toggle' out "
        "loud.",
    ),
    (
        5,
        "Demo: memory off, then memory on.",
        [
            "Off — 'I'm allergic to shellfish, and I'm cooking for six on Friday.' "
            "New thread. 'What should I make Friday?' It suggests shrimp.",
            "On — same two messages, verbatim. New thread again. 'For six, and no "
            "shellfish — here's a menu.' It never re-asked.",
            "The new thread is the proof: nothing was in the context window.",
            "Then switch the user ID. The other user gets nothing — isolation is "
            "enforced in the query, not in the prompt.",
        ],
        "2:00. The 'new thread' step is the whole demo — say it out loud: 'new "
        "conversation, nothing in the context window.' Same thread would be "
        "indistinguishable from ordinary chat history. The user switch is the "
        "second-most convincing five seconds in the talk. FALLBACK: if the recording "
        "will not play, describe it in twenty seconds and jump to slide 6, which is a "
        "still that carries the same point. Never debug live at a booth — five second "
        "rule, then move.",
    ),
    (
        6,
        "What it remembered, and why that memory won.",
        [
            # No literal scores on this slide. The screenshot is the source of truth,
            # and a printed 0.9 beside a screen showing 0.82 is the one error a crowd
            # is guaranteed to catch.
            "Long-term — 'no shellfish', with its importance score and recall count.",
            "Episodic — the turns it logged, with step numbers and the tools that ran.",
            "Short-term — this thread's state, expiring on its own.",
            "Cache — a miss, then a hit, at zero inference cost.",
            "Every line is a document the agent read, shown with the score that got it "
            "there. This is not an abstraction over memory. It is the memory.",
        ],
        "1:30. READ EVERY NUMBER OFF THE SCREEN, never from these notes — whatever the "
        "seeded values are that morning is what you say. Point at the importance score: "
        "it was not hand-written, an LLM assessed it at write time and it climbs every "
        "time the fact gets used. That is how the agent decides what is worth keeping "
        "when there is more history than context. Then point at the cache hit — second "
        "identical question, zero inference cost. Quote the similarity decimal only if "
        "it is legible on screen. Scores must read from eight feet; if they do not, "
        "that is a re-capture, not a live zoom.",
    ),
    (
        7,
        "In Compass: they are just documents.",
        [
            "A long-term memory document — content, importance, access_count, and the "
            "embedding, side by side.",
            "An episodic document — messages, tool calls, files touched, step and "
            "parent_step, under a correlation ID.",
            "The index list — vector, full-text, and TTL indexes, side by side in one "
            "collection list.",
            "Same cluster, same database. No sync job, no ETL, no second system. The "
            "documents the agent retrieves are the documents it writes.",
        ],
        "1:30. Land the TTL point here, because it is what the ops person in the crowd "
        "is already worrying about: short-term state expires on its own — that is an "
        "index, not a cron job you have to remember to write. Retention on the "
        "activity log is one collMod away. DO NOT QUOTE INDEX COUNTS — Compass is "
        "showing the real list next to you, so point at it and say 'vector, full-text, "
        "and TTL, same collection list.' (Reference only, not for speaking: 3 vector, "
        "2 full-text, 7 TTL as the schema stands today.) Every tab pre-opened and "
        "pre-scrolled; do not connect to a cluster in front of the crowd. Static "
        "screenshots are a perfectly good substitute.",
    ),
    (
        8,
        "The retrieval layer, in one aggregation.",
        # Four callouts, in the order they are pointed at. The pipeline itself belongs
        # on the slide for the photo; at booth distance a wall of JSON reads as
        # texture, so the words carry the room and the code carries the camera.
        [
            "1 — $vectorSearch: meaning.",
            "2 — $search: exact terms. SKUs, error codes, names.",
            "3 — $rankFusion: merges both rankings natively, weighted. One round trip.",
            "4 — the filter, inside BOTH branches: per-user isolation, enforced by the "
            "engine rather than by hoping the model behaves.",
            "Final ranking is then calibrated in the application across recency, "
            "importance, and relevance — because 'most similar' is not the same as "
            "'most worth showing'.",
        ],
        "1:30. The technical high point and the slide people photograph — pause three "
        "seconds at the top before you talk. NEVER READ THE CODE. Point at the four "
        "numbered callouts, one clause each: 'Meaning. Exact terms. Merged natively. "
        "Filtered per user, in both branches.' Then talk to the callouts, not the JSON. "
        "Vector alone misses exact terms; full-text alone misses meaning; merging two "
        "ranked lists is the part everyone gets wrong. In the built deck, dim the code "
        "block to about 60 percent and keep the four callouts at full contrast — anyone "
        "who wants the pipeline will photograph it, anyone at the back needs the four "
        "words. If asked: $rankFusion needs MongoDB 8.1+, and before that you compose "
        "it with $unionWith and do the RRF math yourself, which the library supports.",
    ),
    (
        9,
        "Four boxes collapse into one cluster.",
        [
            "Before — a vector database, a cache, a search engine, and a system of "
            "record, stitched together with sync jobs.",
            "After — one Atlas cluster: memories, episodes, semantic_cache, decisions, "
            "audit_log.",
            "$vectorSearch, $search, $rankFusion, and TTL indexes, over the same "
            "documents.",
            "Every arrow you delete is a failure mode you delete. Four systems means "
            "four consistency models, four security reviews, four on-call rotations, "
            "and a sync job that is always a little bit behind.",
        ],
        "1:00. THIS IS THE CUT SLIDE. If you are past 9:45 leaving slide 8, or the "
        "crowd is thinning, SKIP IT ENTIRELY and go to slide 10 — slide 10 is the "
        "differentiator, this one is only the argument, and losing the argument costs "
        "nothing because slide 11 makes it in two lines and Thursday's main session IS "
        "this slide for forty-five minutes. COMPRESSED VERSION (0:20): show both "
        "diagrams, say 'four systems, four consistency models, four on-call rotations — "
        "or one cluster', advance. No trade-off discussion, no sync-job riff. FULL "
        "VERSION: the word to use is 'collapse', because it hands off to Thursday. Do "
        "not oversell it as free — if asked, the trade is real: best-of-breed "
        "specialization for one operational surface. For agent memory specifically that "
        "trade is easy, because the workload is small vectors, a high write rate, and "
        "mandatory per-user filtering, which is a document-database shape.",
    ),
    (
        10,
        "The one people skip: episodic memory.",
        [
            "Semantic memory knows facts. Episodic memory knows what happened.",
            "One document per turn: what was said, which tools ran, what was planned, "
            "which files were touched, in what order, under which trace ID.",
            "Ask your agent 'what did we decide about the Q3 forecast, and what did you "
            "actually change?' — semantic memory cannot answer that.",
            "It has to be free at the call site: fire-and-forget onto a bounded queue, "
            "batched to Atlas by a background worker. Queue full drops the oldest turn, "
            "never the newest. A failed write is a counter, not an exception.",
            "Append-only, TTL'd, hybrid-searchable per user, and never on the hot path. "
            "Your agent never waits on its own diary.",
        ],
        "1:15. This is the slide that separates this talk from every other memory talk "
        "at the show — protect it when time is tight. If asked why not use an "
        "observability vendor: those are built for humans debugging after the fact; "
        "this is built for the agent to query mid-loop, per user, with the same hybrid "
        "search as everything else, next to the rest of your data.",
    ),
    (
        11,
        "Context engineering is a data problem. Solve it in the data layer.",
        [
            "Four memories, one cluster: state, facts, events, cache.",
            "Less repetition, and a product that compounds in value the longer someone "
            "uses it.",
            "More consistent behavior across sessions, with less prompting.",
            "Lower token spend than replaying transcripts, and fewer moving parts than "
            "a split architecture.",
        ],
        "0:30. Say the two title lines, then stop. Let it sit for a beat — this is the "
        "sentence you want them repeating at the next booth.",
    ),
    (
        12,
        "Build agents that remember on purpose.",
        # Ordered cheapest commitment first: install, repo, blog, Thursday. The order
        # is the ask.
        [
            # Not on PyPI. A `pip install agent-memory` on screen resolves to somebody
            # else's package or nothing at all, and someone in the front row will try it.
            "1 — Install it:  uv add git+https://github.com/mongodb-partners/ai-memory.git",
            "2 — Read it:  github.com/mongodb-partners/ai-memory — framework-neutral "
            "Python library, callable directly, over MCP, or over REST. Sample UI with "
            "the memory panel and the Compass pipelines included.",
            "3 — Go deeper:  the blog, 'Build AI Memory Systems with MongoDB Atlas, "
            "AWS & Claude'.",
            "4 — See the whole story:  my main session, Thursday, August 6th — 'For "
            "Agents: Deployment at Scale'.",
        ],
        "0:15. Four things, in this order, in one breath — the order is the ask, from "
        "cheapest commitment to biggest. 'INSTALL IT STRAIGHT FROM THE REPO — one "
        "command, it is open source. REPO is on screen, sample UI is in it. BLOG if you "
        "want the long "
        "version. And THE WHOLE STACK-COLLAPSE STORY is my main session Thursday — come "
        "to that one.' QR to the repo goes on the booth backdrop, not the slide: the "
        "slide is up for fifteen seconds and nobody scans a moving target. Then stop "
        "talking and take one question. Do not start a second talk.",
    ),
]


def set_body(shape, paragraphs: list[str]) -> None:
    """Replace a body text box's paragraphs, cloning the existing paragraph style.

    Every body paragraph in the base deck carries an identical ``<a:pPr>`` (font size,
    colour, space-after). We deep-copy the first paragraph's XML per line so the new
    text inherits it exactly, instead of guessing at the formatting.
    """
    tf = shape.text_frame
    template = copy.deepcopy(tf.paragraphs[0]._p)
    body = tf._txBody
    for p in list(body.findall(f"{{{tf.paragraphs[0]._p.nsmap['a']}}}p")):
        body.remove(p)

    for text in paragraphs:
        p = copy.deepcopy(template)
        ns = f"{{{template.nsmap['a']}}}"
        # Drop every run but the first, then set that run's text.
        runs = p.findall(f"{ns}r")
        for extra in runs[1:]:
            p.remove(extra)
        if runs:
            t = runs[0].find(f"{ns}t")
            t.text = text
        body.append(p)


def set_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def by_name(slide, name: str):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    raise KeyError(f"no shape named {name!r} on this slide")


def build(base: Path, out: Path) -> None:
    prs = Presentation(str(base))
    if len(prs.slides) < 12:
        sys.exit(f"expected at least 12 slides in {base}, found {len(prs.slides)}")

    # Slide 1: subtitle + notes only. Its title, date block, and pull-quote are correct.
    title_slide = prs.slides[0]
    by_name(title_slide, TITLE_SLIDE_SUBTITLE).text_frame.paragraphs[0].runs[
        0
    ].text = TITLE_SLIDE_SUBTITLE_TEXT
    set_notes(title_slide, TITLE_SLIDE_NOTES)

    for number, title, paragraphs, notes in SLIDES:
        slide = prs.slides[number - 1]
        title_box = by_name(slide, TITLE_SHAPE)
        title_box.text_frame.paragraphs[0].runs[0].text = title
        set_body(by_name(slide, BODY_SHAPE), paragraphs)
        set_notes(slide, notes)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"wrote {out}")
    print(f"  {len(prs.slides)} slides, speaker notes on all 12")


def main() -> None:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--base", type=Path, default=DEFAULT_BASE, help="styled base deck")
    ap.add_argument("--out", type=Path, default=DEFAULT_OUT, help="output .pptx")
    args = ap.parse_args()

    if not args.base.exists():
        sys.exit(f"base deck not found: {args.base}")
    build(args.base, args.out)


if __name__ == "__main__":
    main()
